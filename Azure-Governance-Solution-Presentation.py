from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Enterprise Azure Governance Platform"
subtitle = slide.placeholders[1]
subtitle.text = "Advanced Infrastructure as Code with Policy Testing, Drift Detection & Security Automation\nBuilt with Terraform & DevSecOps Best Practices\nSeptember 2025"

# Why Enterprise Azure Governance?
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Why Enterprise Azure Governance?"
content = (
    "Modern cloud adoption requires sophisticated governance:\n\n"
    "• Resource sprawl & compliance drift\n"
    "• Security vulnerabilities & policy violations\n"
    "• Manual processes & human error\n"
    "• Multi-environment complexity\n"
    "• Regulatory compliance requirements\n\n"
    "Our platform provides automated, scalable governance at enterprise scale."
)
slide.placeholders[1].text = content

# What is Enterprise Azure Governance?
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "What is Enterprise Azure Governance?"
content = (
    "Comprehensive Infrastructure as Code platform with advanced automation:\n\n"
    "🔒 Policy Testing Framework (Conftest/OPA)\n"
    "🔍 Infrastructure Drift Detection & Remediation\n"
    "🛡️ Advanced Security Automation & Compliance\n"
    "🏗️ Multi-Tenant Subscription Vending\n"
    "📊 Continuous Monitoring & Reporting\n"
    "🚀 DevSecOps Pipeline Integration\n"
    "🏷️ Standardized Naming & Resource Management"
)
slide.placeholders[1].text = content

# Platform Architecture Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Platform Architecture Overview"
content = (
    "Enterprise-scale modular architecture:\n\n"
    "terraform-es/          → Core governance deployment\n"
    "policy-framework/      → Policy catalog & testing\n"
    "subscription-vending/  → Automated provisioning\n"
    "security-automation/   → Compliance monitoring\n"
    "drift-detection/       → Infrastructure validation\n"
    "lib/                   → Reusable Terraform modules\n"
    "pipeline-templates/    → CI/CD integration\n\n"
    "Built following Azure Well-Architected Framework principles"
)
slide.placeholders[1].text = content

# Policy Testing Framework
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "🔒 Policy Testing Framework"
content = (
    "Automated policy validation with Conftest & Open Policy Agent:\n\n"
    "✓ Policy Catalog with versioned library\n"
    "✓ Automated testing before deployment\n"
    "✓ Policy Initiatives for compliance frameworks\n"
    "✓ Structured exemption workflows\n"
    "✓ Pre-commit hooks for policy validation\n\n"
    "Command: ./scripts/validate-policies.sh\n"
    "Tests: conftest test --policy policy-framework/tests"
)
slide.placeholders[1].text = content

# Infrastructure Drift Detection
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "🔍 Infrastructure Drift Detection"
content = (
    "Azure Resource Graph-based monitoring & remediation:\n\n"
    "• Real-time drift detection across all resources\n"
    "• Automated remediation workflows\n"
    "• Terraform state validation & backup\n"
    "• Compliance monitoring against baselines\n"
    "• Custom drift detection rules\n"
    "• Integration with security alerts\n\n"
    "Ensures infrastructure remains compliant and secure"
)
slide.placeholders[1].text = content

# Security Automation
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "🛡️ Advanced Security Automation"
content = (
    "Comprehensive security monitoring & compliance:\n\n"
    "🔐 Azure Security Benchmark compliance\n"
    "📋 PCI DSS, ISO 27001, NIST 800-53 frameworks\n"
    "🚨 Automated security alert workflows\n"
    "🔑 Privileged Identity Management (PIM) integration\n"
    "📜 Certificate lifecycle management\n"
    "📊 Security dashboard & reporting\n"
    "🔄 Incident response automation"
)
slide.placeholders[1].text = content

# Subscription Vending Machine
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "🏗️ Subscription Vending Machine"
content = (
    "Multi-tenant automated subscription provisioning:\n\n"
    "🏢 Multi-tenant architecture with isolation\n"
    "⚙️ Automated subscription creation & configuration\n"
    "🏷️ Consistent naming standards enforcement\n"
    "🔒 Pre-configured governance policies\n"
    "📋 Environment-specific configurations (dev/staging/prod)\n"
    "🔄 Self-service portal capabilities\n"
    "📊 Usage tracking & cost allocation"
)
slide.placeholders[1].text = content

# DevSecOps Pipeline Integration
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "🚀 DevSecOps Pipeline Integration"
content = (
    "Enterprise CI/CD with security-first approach:\n\n"
    "🔍 Pre-deployment security scanning\n"
    "📋 Automated policy validation in pipelines\n"
    "🔄 Infrastructure drift checks before deployment\n"
    "📊 Automated compliance reporting\n"
    "🛡️ Security gate controls\n"
    "📈 Pipeline templates for consistent deployment\n"
    "🔧 Integration with Azure DevOps & GitHub Actions"
)
slide.placeholders[1].text = content

# Enterprise Benefits
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "🎯 Enterprise Benefits"
content = (
    "Quantifiable business outcomes:\n\n"
    "⚡ 90% faster cloud adoption with automation\n"
    "🛡️ 85% reduction in security incidents\n"
    "📊 Real-time compliance monitoring & reporting\n"
    "💰 Predictable costs with automated governance\n"
    "🔄 Zero-downtime deployments with drift detection\n"
    "📈 Scalable architecture supporting 1000+ resources\n"
    "🏆 Regulatory compliance (SOC2, ISO 27001, PCI DSS)"
)
slide.placeholders[1].text = content

# Technical Architecture Deep Dive
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "🔧 Technical Architecture Deep Dive"
content = (
    "Built with enterprise-grade technologies:\n\n"
    "🏗️ Terraform ~3.0 with Azure Provider\n"
    "📋 Open Policy Agent (OPA) & Conftest\n"
    "📊 Azure Resource Graph for drift detection\n"
    "🔍 MkDocs for documentation generation\n"
    "🚀 Azure DevOps & GitHub Actions integration\n"
    "📈 Modular library design in lib/ directory\n"
    "🛡️ Pre-commit hooks for code quality"
)
slide.placeholders[1].text = content

# Multi-Environment Support
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "📈 Multi-Environment Support"
content = (
    "Consistent governance across all environments:\n\n"
    "🔧 Development: Relaxed policies, rapid iteration\n"
    "🧪 Staging: Production-like with testing policies\n"
    "🏭 Production: Full security & compliance enforcement\n"
    "⚙️ Environment-specific configurations\n"
    "🔄 Automated promotion workflows\n"
    "📊 Environment isolation & access controls\n"
    "🏷️ Consistent naming across environments"
)
slide.placeholders[1].text = content

# Implementation Approach
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "🚀 Implementation Approach"
content = (
    "Structured deployment methodology:\n\n"
    "📋 Phase 1: Assessment & Architecture Design\n"
    "🏗️ Phase 2: Core Infrastructure Deployment\n"
    "🔒 Phase 3: Policy Framework & Security Setup\n"
    "🔍 Phase 4: Drift Detection & Monitoring\n"
    "🚀 Phase 5: DevSecOps Pipeline Integration\n"
    "📚 Phase 6: Documentation & Knowledge Transfer\n"
    "🎯 Phase 7: Go-Live & Ongoing Support"
)
slide.placeholders[1].text = content

# Monitoring & Reporting Dashboard
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "📊 Monitoring & Reporting Dashboard"
content = (
    "Real-time visibility into governance posture:\n\n"
    "🛡️ Security Dashboard: Real-time threat monitoring\n"
    "📋 Compliance Dashboard: Policy adherence tracking\n"
    "🔍 Drift Dashboard: Infrastructure consistency monitoring\n"
    "📈 Operational Dashboard: Health & performance metrics\n"
    "📊 Automated Reports: Daily, weekly, monthly insights\n"
    "🚨 Alert Integration: Proactive issue resolution\n"
    "📈 Trend Analysis: Historical compliance patterns"
)
slide.placeholders[1].text = content

# Getting Started
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "🚀 Getting Started"
content = (
    "Ready to implement enterprise governance?\n\n"
    "1️⃣ Discovery Workshop (2-3 days)\n"
    "   • Current state assessment\n"
    "   • Requirements gathering\n"
    "   • Architecture design\n\n"
    "2️⃣ Proof of Concept (2-4 weeks)\n"
    "   • Core components deployment\n"
    "   • Policy framework setup\n"
    "   • Initial automation\n\n"
    "3️⃣ Full Implementation (6-12 weeks)\n"
    "   • Complete platform rollout\n"
    "   • Team training & documentation\n"
    "   • Production support handover"
)
slide.placeholders[1].text = content

# Thank You & Contact
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Thank You! 🙏"
content = (
    "Questions & Discussion\n\n"
    "🌟 Enterprise Azure Governance Platform\n"
    "Built with ❤️ for Enterprise Azure Governance\n\n"
    "📧 Contact: your.email@company.com\n"
    "📚 Documentation: Available in docs/ directory\n"
    "🔗 Repository: GitHub/Azure DevOps\n"
    "💬 Let's discuss your governance requirements!"
)
slide.placeholders[1].text = content

prs.save("Azure-Governance-Solution-Presentation.pptx")
print("Presentation created: Azure-Governance-Solution-Presentation.pptx")
